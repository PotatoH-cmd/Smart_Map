# -*- coding: utf-8 -*-
"""生成《豫水智能一张图》项目总结 PPT（16:9）。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 配色 ──
DEEP  = RGBColor(0x0F, 0x3D, 0x6E)   # 深蓝
MID   = RGBColor(0x2E, 0x8B, 0xC0)   # 中蓝
CYAN  = RGBColor(0x27, 0xA9, 0xE0)   # 亮青
LIGHT = RGBColor(0xE8, 0xF3, 0xFB)   # 浅蓝底
PALE  = RGBColor(0xF5, 0xF9, 0xFD)   # 极浅底
INK   = RGBColor(0x22, 0x30, 0x3E)   # 正文深色
GRAY  = RGBColor(0x64, 0x74, 0x8B)   # 次要文字
AMBER = RGBColor(0xF0, 0x9A, 0x1F)   # 强调橙
GREEN = RGBColor(0x2E, 0x9E, 0x6B)   # 强调绿
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE  = RGBColor(0xD5, 0xE3, 0xF0)   # 描边

FONT = "微软雅黑"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── 基础工具 ──
def set_cjk(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', name)

def box(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp

def txt(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=0, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (text, size, bold, color)"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after: p.space_after = Pt(space_after)
        if line_spacing: p.line_spacing = line_spacing
        for (t, s, b, c) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(s); r.font.bold = b; r.font.color.rgb = c
            set_cjk(r)
    return tb

def card(slide, x, y, w, h, title, body, fill=PALE, accent=MID, title_size=15, body_size=11, title_color=None, pad=0.14):
    """带顶部色条的卡片"""
    box(slide, x, y, w, h, fill=fill, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    box(slide, x, y, 0.09, h, fill=accent)
    txt(slide, x + pad + 0.06, y + 0.12, w - pad * 2 - 0.1, 0.35,
        [[(title, title_size, True, title_color or DEEP)]])
    txt(slide, x + pad + 0.06, y + 0.5, w - pad * 2 - 0.1, h - 0.6,
        [[(ln, body_size, False, INK)] for ln in body], line_spacing=1.12, space_after=3)

def chip(slide, x, y, w, h, text, fill, color=WHITE, size=12, bold=True):
    sp = box(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = sp.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    set_cjk(r)
    return sp

def arrow(slide, x, y, w=0.45, h=0.3, color=MID, direction="right"):
    shp = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    sp = box(slide, x, y, w, h, fill=color, shape=shp)
    return sp

PAGE_NO = [0]
def content_slide(title, tag=None):
    s = prs.slides.add_slide(BLANK)
    box(s, 0, 0, 13.333, 7.5, fill=WHITE)
    # 顶部装饰条
    box(s, 0, 0, 13.333, 0.09, fill=DEEP)
    box(s, 0.5, 0.46, 0.13, 0.52, fill=CYAN)
    txt(s, 0.78, 0.42, 9.5, 0.6, [[(title, 27, True, DEEP)]])
    if tag:
        chip(s, 11.15, 0.5, 1.7, 0.42, tag, LIGHT, color=MID, size=12)
    txt(s, 9.9, 0.52, 1.2, 0.4, [[("Map Assistant", 10, False, GRAY)]], align=PP_ALIGN.RIGHT)
    # 页脚
    box(s, 0.5, 7.06, 12.33, 0.012, fill=LINE)
    PAGE_NO[0] += 1
    txt(s, 0.5, 7.12, 8, 0.3, [[("豫水智能一张图 · Map Assistant v1", 9, False, GRAY)]])
    txt(s, 12.3, 7.12, 0.55, 0.3, [[(str(PAGE_NO[0]).zfill(2), 9, True, GRAY)]], align=PP_ALIGN.RIGHT)
    return s

# ══════════════════════ 第 1 页 · 封面 ══════════════════════
s = prs.slides.add_slide(BLANK)
box(s, 0, 0, 13.333, 7.5, fill=DEEP)
box(s, 9.2, -1.6, 5.4, 5.4, fill=MID, shape=MSO_SHAPE.OVAL)
box(s, 10.7, -0.4, 2.9, 2.9, fill=CYAN, shape=MSO_SHAPE.OVAL)
box(s, -1.8, 5.6, 4.4, 4.4, fill=MID, shape=MSO_SHAPE.OVAL)
box(s, 0.5, 1.15, 0.9, 0.14, fill=CYAN)
txt(s, 0.5, 1.55, 11, 0.5, [[("豫水智能一张图", 17, False, RGBColor(0xBF, 0xD9, 0xF2))]])
txt(s, 0.5, 2.55, 12.3, 1.2, [[("豫水智能一张图", 52, True, WHITE)]])
txt(s, 0.5, 3.75, 12.3, 0.7, [[("AI 驱动的地理空间智能分析平台 · Map Assistant v1", 24, False, CYAN)]])
txt(s, 0.5, 4.85, 12.3, 0.5,
    [[("面向水利勘测  ·  河道采砂监管  ·  河道侵占识别", 15, False, RGBColor(0xCF, 0xE4, 0xF6))]])
txt(s, 0.5, 6.35, 12.3, 0.5,
    [[("项目总结汇报材料", 13, False, RGBColor(0x9F, 0xC3, 0xE4))]])
txt(s, 9.0, 6.35, 3.9, 0.5, [[("2026", 13, False, RGBColor(0x9F, 0xC3, 0xE4))]], align=PP_ALIGN.RIGHT)

# ══════════════════════ 第 2 页 · 目录 ══════════════════════
s = content_slide("目录", "CONTENTS")
items = [
    ("01", "建设背景与业务痛点", "GIS 门槛高 · 巡查效率低 · 遥感周期长"),
    ("02", "项目定位与核心能力", "对话即操作，九大核心能力全景"),
    ("03", "总体技术架构", "数据 / AI 引擎 / 服务 / 展示 四层协同"),
    ("04", "智能体编排与工具体系", "IntentAgent → Harness → LangGraph · 15+ 工具"),
    ("05", "遥感检测与智能报告", "SAM3 三阶段流水线 · 图斑后处理 · Word 报告"),
    ("06", "应用场景与未来展望", "河道监管 / 采砂管理 / 决策报表 / 培训问答"),
]
for i, (num, t, d) in enumerate(items):
    col = i % 2; row = i // 2
    x = 0.8 + col * 6.1; y = 1.55 + row * 1.7
    box(s, x, y, 5.75, 1.35, fill=PALE, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    box(s, x + 0.22, y + 0.28, 0.8, 0.8, fill=DEEP if col == 0 else MID, shape=MSO_SHAPE.OVAL)
    txt(s, x + 0.22, y + 0.42, 0.8, 0.5, [[(num, 18, True, WHITE)]], align=PP_ALIGN.CENTER)
    txt(s, x + 1.25, y + 0.2, 4.4, 0.45, [[(t, 16, True, DEEP)]])
    txt(s, x + 1.25, y + 0.72, 4.4, 0.45, [[(d, 10.5, False, GRAY)]])

# ══════════════════════ 第 3 页 · 痛点 ══════════════════════
s = content_slide("建设背景与业务痛点", "01")
pains = [
    ("GIS 软件门槛高", "专业空间分析依赖 ArcGIS / QGIS 等专业软件，基层人员难以掌握"),
    ("河道巡查效率低", "人工比对影像与红线，侵占行为发现周期长、覆盖不全"),
    ("遥感分析周期长", "影像分割、变化检测、图斑筛选流程繁琐，单次分析数天"),
    ("报表工作重复枯燥", "周报 / 月报需人工拼接数据、截图、图表，重复劳动大"),
    ("政策法规查询难", "规范条文分散于大量文档，检索费时、引用易出错"),
]
for i, (t, d) in enumerate(pains):
    x = 0.5 + i * 2.5
    card(s, x, 1.5, 2.35, 3.3, t, [d], accent=DEEP if i % 2 == 0 else CYAN, body_size=11)
box(s, 0.5, 5.35, 12.33, 0.95, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
txt(s, 0.9, 5.55, 11.6, 0.6, [
    [("目标：", 15, True, AMBER), ("用 AI 降低专业工具使用门槛 —— 一句话完成空间分析，让科技服务基层水利管理", 15, True, DEEP)],
], align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════ 第 4 页 · 定位与核心能力 ══════════════════════
s = content_slide("项目定位与核心能力全景", "02")
box(s, 0.5, 1.35, 12.33, 1.15, fill=DEEP, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
txt(s, 0.9, 1.55, 11.6, 0.8, [
    [("一句话定位：", 15, True, CYAN), ("面向水利行业的多模态 AI 地理空间分析平台 —— 将大语言模型、计算机视觉与 GIS 深度融合，", 15, True, WHITE)],
    [("以自然语言对话为核心入口，让专业 GIS 分析“像和人对话一样简单”。", 15, True, WHITE)],
], line_spacing=1.15)
caps = [
    ("智能对话", "LangGraph 多 Agent 协作\n意图自动识别与分派"),
    ("2D/3D 地图", "Leaflet + Cesium 双引擎\n矢量图层 / 高分影像叠加"),
    ("知识库问答", "水利规范 / 政策检索\nRagFlow / LlamaIndex 双后端"),
    ("空间分析", "缓冲区 / 裁剪 / 叠加\nQGIS MCP 引擎驱动"),
    ("数据库查询", "自然语言转 SQL\n直连 PostgreSQL + PostGIS"),
    ("采砂报告生成", "RTK 测点 + 无人机影像\n一键输出 Word 报告"),
    ("遥感变化检测", "SAM3 分割 + 双时相检测\n河道侵占自动识别"),
    ("切片管理", "MVT / 3D Tiles / MBTiles\n发布与异步构建"),
    ("图表可视化", "对话式柱状 / 饼图 / 折线\nECharts 实时渲染"),
]
for i, (t, d) in enumerate(caps):
    x = 0.5 + (i % 3) * 4.15; y = 2.75 + (i // 3) * 1.4
    card(s, x, y, 3.95, 1.25, t, d.split("\n"), accent=MID, title_size=13, body_size=10, pad=0.1)

# ══════════════════════ 第 5 页 · 总体架构 ══════════════════════
s = content_slide("总体技术架构 —— 四层协同", "03")
layers = [
    ("展示层", CYAN, ["React 18 SPA：Leaflet 2D + Cesium 3D 双地图引擎", "ECharts 图表 · 聊天面板 · 工具面板 · Word 报告在线预览"]),
    ("服务层", MID, ["FastAPI（:8006）· 88+ API 端点 · SSE 流式对话 / WebSocket", "会话管理 · 文件上传 · 切片服务 · GeoServer REST 代理 · PM2 托管"]),
    ("AI 引擎层", DEEP, ["LangGraph 状态图 + Qwen-Agent 工具基座（IntentAgent → Harness → TaskExecutor）", "千问 LLM · Ollama 视觉模型 · SAM3 · RAGFlow / LlamaIndex"]),
    ("数据层", RGBColor(0x5B, 0x3A, 0x8E), ["PostgreSQL + PostGIS（业务 / 空间数据）· SQLite（会话持久化）", "GeoServer WMS/WFS · 高分影像 · 无人机航拍 · MVT 矢量瓦片 · 3D Tiles"]),
]
y = 1.35
for i, (name, color, lines) in enumerate(layers):
    box(s, 1.3, y, 10.7, 1.05, fill=PALE, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    chip(s, 1.55, y + 0.22, 1.75, 0.6, name, color, size=14)
    txt(s, 3.55, y + 0.16, 8.2, 0.8, [[(ln, 11, False, INK)] for ln in lines], line_spacing=1.15)
    if i < 3:
        arrow(s, 6.45, y + 1.07, 0.4, 0.24, color=MID, direction="down")
    y += 1.4
txt(s, 0.5, 6.75, 12.3, 0.3, [[("外部依赖：PostgreSQL + PostGIS · GeoServer · QGIS MCP (Docker) · RagFlow / LlamaIndex · DashScope API", 10.5, False, GRAY)]])

# ══════════════════════ 第 6 页 · 智能体编排 ══════════════════════
s = content_slide("AI 智能对话 —— 从“一句话”到“全自动分析”", "04")
flow = [
    ("用户输入", "自然语言\n一句话描述需求"),
    ("意图识别", "IntentAgent\nLLM 结构化意图\n+ 60+ 快速关键词路由"),
    ("Agent 分派", "AgentHarness\n意图 → 领域 Agent\n提示词 / 上下文组装"),
    ("工具执行", "TaskExecutor\nLangGraph 三节点图\n实时事件推送"),
    ("结果合成", "Response\n自然语言回复\n+ 地图 / 图表指令"),
]
x = 0.5
for i, (t, d) in enumerate(flow):
    w = 2.15
    box(s, x, 1.5, w, 1.5, fill=LIGHT if i != 3 else RGBColor(0xDD, 0xEE, 0xFA), line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    txt(s, x + 0.1, 1.62, w - 0.2, 0.4, [[(t, 14, True, DEEP)]], align=PP_ALIGN.CENTER)
    txt(s, x + 0.12, 2.05, w - 0.24, 0.9, [[(ln, 9.5, False, GRAY)] for ln in d.split("\n")], align=PP_ALIGN.CENTER, line_spacing=1.1)
    if i < 4:
        arrow(s, x + w + 0.04, 2.05, 0.31, 0.28, color=CYAN)
    x += w + 0.35
txt(s, 0.5, 3.35, 12.3, 0.35, [[("五大领域 Agent：按意图自动分派，可插拔扩展", 12.5, True, MID)]])
agents = [
    ("MapAgent", "地图展示 · 位置搜索\n空间处理 / 空间参考"),
    ("DataAgent", "数据查询\n可视化"),
    ("KnowledgeAgent", "知识库检索\n规范问答"),
    ("ReportAgent", "报告生成\n采砂报告"),
    ("GeneralAgent", "兜底 · 跨意图复合\n天气 / 未知意图"),
]
for i, (t, d) in enumerate(agents):
    x = 0.5 + i * 2.5
    card(s, x, 3.8, 2.35, 1.5, t, d.split("\n"), accent=DEEP if i % 2 == 0 else CYAN, title_size=13, body_size=10)
box(s, 0.5, 5.6, 12.33, 1.1, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
txt(s, 0.9, 5.78, 11.6, 0.8, [
    [("执行示例：", 13, True, AMBER), ("“查看潢河郝楼可采区的河道侵占情况”", 13, True, DEEP)],
    [("空间关键词识别 → 获取红线边界 → 查询侵占图斑 → 加载地图 → 生成分析文档（全程 LangGraph 自动编排）", 12, False, INK)],
], line_spacing=1.2)

# ══════════════════════ 第 7 页 · 工具体系 ══════════════════════
s = content_slide("工具体系 —— 15+ 模块协同工作", "04")
groups = [
    ("地图与 3D 场景", ["map_tool — 2D 地图操作", "cesium_tool — 3D 飞行 / 叠加", "overlay / satellite 瓦片服务"]),
    ("数据查询", ["postgresql_tool — NL2SQL", "mcp_postgres_tool — MCP 查询", "schema_manager — 表结构管理"]),
    ("空间分析", ["qgis_mcp_tool — QGIS 引擎", "spatial_processing — 坐标转换", "spatial_reference — 红线/采区"]),
    ("遥感与 AI", ["sam_predict — SAM3 分割", "resam / ssa — 检测模型训练", "weather_tool — 天气查询"]),
    ("知识库", ["ragflow_knowledge_tool", "llamaindex_knowledge_tool", "knowledge_graph — Kuzu 图谱"]),
    ("报告与可视化", ["report_generator_tool — Word", "caisha_report_tool — 采砂报告", "data_visualizer — ECharts 图表"]),
]
for i, (t, lines) in enumerate(groups):
    x = 0.5 + (i % 3) * 4.15; y = 1.45 + (i // 3) * 2.6
    card(s, x, y, 3.95, 2.4, t, lines, accent=MID, title_size=14, body_size=11, pad=0.12)
box(s, 0.5, 6.72, 12.33, 0.0, fill=None)

# ══════════════════════ 第 8 页 · 遥感检测 ══════════════════════
s = content_slide("多模态遥感检测 —— 河道侵占自动识别", "05")
stages = [
    ("① SAM3 全图分割", ["高分影像地物分割", "水体 / 建筑 / 植被 / 道路", "双头掩码融合提升小目标精度", "前端 SAMPanel 交互标注"]),
    ("② 双时相变化检测", ["OmniOVCD 变化检测框架", "自动发现新增建筑 / 道路扩张", "本地 Ollama 视觉模型", "DashScope 千问 VL 双后端"]),
    ("③ 图斑智能后处理", ["重叠合并 · 聚集密度分析", "大面积单体规则（≥5000㎡）", "VLM 二次确认", "归类：确认 / 疑似 / 排除"]),
]
x = 0.5
for i, (t, lines) in enumerate(stages):
    w = 3.75
    card(s, x, 1.5, w, 3.6, t, lines, accent=DEEP if i == 0 else (MID if i == 1 else AMBER), title_size=15, body_size=11.5)
    if i < 2:
        arrow(s, x + w + 0.1, 3.1, 0.42, 0.3, color=CYAN)
    x += w + 0.42
box(s, 0.5, 5.5, 12.33, 1.05, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
txt(s, 0.9, 5.66, 11.6, 0.75, [
    [("三阶段流水线全自动运行：", 13.5, True, DEEP), ("影像 → 图斑 → 结论，无需人工干预；检测结果可直接叠加到 2D/3D 地图并生成报告", 13.5, False, INK)],
    [("试点：潢河郝楼、白露河黄堰等采区 · 已处理数十平方公里高分影像 · 识别数百处疑似侵占图斑", 11.5, False, GRAY)],
], line_spacing=1.2)

# ══════════════════════ 第 9 页 · 图斑后处理 ══════════════════════
s = content_slide("图斑智能后处理 —— 从碎片到可操作结论", "05")
flow = [
    ("原始检测图斑", "4,800+ 个碎片"),
    ("重叠合并", "buffer-union-shrink\n合并为 ~2,900 个"),
    ("聚集密度分析", "聚集度评分\n空间密度评估"),
    ("规则判定", "大面积单体 ≥5000㎡\n自动标记可疑"),
    ("VLM 二次确认", "视觉语言模型\n人工复核兜底"),
    ("分类输出", "确认 / 疑似 / 排除\n人工审核量 ↓70%"),
]
x = 0.5
for i, (t, d) in enumerate(flow):
    w = 1.8
    box(s, x, 1.8, w, 1.6, fill=LIGHT if i % 2 == 0 else RGBColor(0xE4, 0xF2, 0xEA), line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    txt(s, x + 0.08, 1.92, w - 0.16, 0.4, [[(t, 12.5, True, DEEP)]], align=PP_ALIGN.CENTER)
    txt(s, x + 0.08, 2.35, w - 0.16, 0.95, [[(ln, 9, False, GRAY)] for ln in d.split("\n")], align=PP_ALIGN.CENTER, line_spacing=1.1)
    if i < 5:
        arrow(s, x + w + 0.02, 2.42, 0.2, 0.24, color=CYAN)
    x += w + 0.24
box(s, 0.5, 4.35, 12.33, 1.5, fill=PALE, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
txt(s, 0.9, 4.55, 11.6, 1.1, [
    [("效果：", 14, True, AMBER), ("原始数千个碎片图斑自动归纳为三类结论，人工审核量降低约 70%", 14, True, DEEP)],
    [("独创“聚集密度分析”：图斑重叠合并 + 聚集密度评估 + 大面积单体规则 + VLM 确认四环节，构建完整自动化质检管道", 12, False, INK)],
], line_spacing=1.25)
txt(s, 0.5, 6.15, 12.3, 0.4, [[("价值：将遥感检测从“研究工具”变成“生产工具”，一线人员直接获得可处置的结论清单", 12, False, GREEN)]])

# ══════════════════════ 第 10 页 · 双地图引擎 ══════════════════════
s = content_slide("双地图引擎 —— 2D 与 3D 无缝协同", "03")
card(s, 0.5, 1.45, 6.0, 3.6, "2D 视图 · Leaflet（日常巡查 / 数据浏览）",
     ["ArcGIS 切片 / 高分影像 / OSM 多底图切换", "图层控制面板：采区边界、河道红线、检测图斑", "点位抽稀算法 + 视口自适应渲染，海量数据流畅", "交互式标注 / 测量 / 截图上传"], accent=MID, body_size=12)
card(s, 6.83, 1.45, 6.0, 3.6, "3D 视图 · CesiumJS（地形分析 / 工程展示）",
     ["3D Tiles 加载 + 无人机航拍模型叠加", "飞行漫游、视角切换、地形分析", "WebSocket 实时指令同步（/ws/cesium）", "三维场景与后端服务双向通信"], accent=DEEP, body_size=12)
box(s, 0.5, 5.45, 12.33, 1.15, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
txt(s, 0.9, 5.62, 11.6, 0.8, [
    [("协同机制：", 14, True, AMBER), ("共享图层数据与空间分析结果，AI 根据上下文自动路由 —— 宏观分析用 2D，地形展示用 3D，操作无缝同步", 14, False, DEEP)],
    [("任一视图中的操作自动同步到另一视图，跨引擎上下文感知（2D 视图优先 Leaflet 工具，3D 视图路由 Cesium 工具）", 11.5, False, GRAY)],
], line_spacing=1.2)

# ══════════════════════ 第 11 页 · 智能报告 ══════════════════════
s = content_slide("智能报告生成 —— 从数据到专业文档", "05")
steps = [
    ("对话分析", "用户一句话描述\n报告需求"),
    ("数据整合", "RTK 测点\n无人机影像\n监测数据"),
    ("地图截图", "前端截图代理\n捕获当前视图\n图层/标注完整"),
    ("知识增强", "自动检索\n水利标准条款\n作为报告依据"),
    ("Word 生成", "模板占位符填充\n图表嵌入\n一键成稿"),
    ("在线预览", "Mammoth 转 HTML\n新标签预览\n打印 / 下载"),
]
x = 0.5
for i, (t, d) in enumerate(steps):
    w = 1.8
    box(s, x, 1.5, w, 1.75, fill=LIGHT if i % 2 == 0 else RGBColor(0xFA, 0xF3, 0xE4), line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    txt(s, x + 0.08, 1.62, w - 0.16, 0.4, [[(t, 13, True, DEEP)]], align=PP_ALIGN.CENTER)
    txt(s, x + 0.08, 2.08, w - 0.16, 1.1, [[(ln, 9.5, False, GRAY)] for ln in d.split("\n")], align=PP_ALIGN.CENTER, line_spacing=1.12)
    if i < 5:
        arrow(s, x + w + 0.02, 2.18, 0.2, 0.24, color=CYAN)
    x += w + 0.24
box(s, 0.5, 3.7, 12.33, 1.35, fill=PALE, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
txt(s, 0.9, 3.88, 11.6, 1.0, [
    [("特色能力：", 13.5, True, AMBER), ("对话即报告 —— 地图截图自动嵌入 + 法规条款自动引用 + 在线预览", 13.5, True, DEEP)],
    [("支持采砂监测专项报告（caisha_report_tool）与通用报告（report_generator_tool），模板由 create_template.py / update_template.py 统一维护", 11.5, False, INK)],
    [("API：/api/preview/report（HTML 预览）· /api/download/report（下载）· /api/save-screenshot（截图上传）", 10.5, False, GRAY)],
], line_spacing=1.25)
txt(s, 0.5, 5.35, 12.3, 0.4, [[("价值：大幅减轻基层周报 / 月报编写负担，让报告既有数据又有空间可视化佐证", 12, False, GREEN)]])

# ══════════════════════ 第 12 页 · 知识库 ══════════════════════
s = content_slide("知识库智能检索 —— 让 AI 回答有据可依", "02")
card(s, 0.5, 1.45, 6.0, 3.5, "双后端架构（KNOWLEDGE_BACKEND 切换）",
     ["RagFlow（默认）：远程知识库服务，文档解析 + 向量检索 + 重排序", "LlamaIndex（本地）：DashScope Embedding + Kuzu 图数据库知识图谱", "一键迁移脚本：migrate_kb_to_llamaindex.py", "知识图谱可视化查询（/api/knowledge/graph）"], accent=MID, body_size=12)
card(s, 6.83, 1.45, 6.0, 3.5, "知识覆盖与智能兜底",
     ["150+ 篇水利标准与法规：采砂规范 / 监测规范 / 实时实施方案", "可视化管理界面：上传 / 更新 / 删除文档（KnowledgeBaseManager）", "智能兜底：数据库查询无结果时自动回退知识库检索", "回答引用原文条款，信息准确可追溯"], accent=DEEP, body_size=12)
box(s, 0.5, 5.3, 12.33, 1.3, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
txt(s, 0.9, 5.48, 11.6, 0.95, [
    [("应用价值：", 13.5, True, AMBER), ("新员工“随身导师” —— 回答河道管理条例、技术标准等专业问题并引用原文；报告生成自动检索法规作为依据", 13.5, False, DEEP)],
    [("导入脚本：ingest_caisha_md.py（采砂规范）· ingest_jiance_md.py（监测规范）· ingest_shishifangan_multi.py（实施方案）", 11, False, GRAY)],
], line_spacing=1.25)

# ══════════════════════ 第 13 页 · 空间智能感知 ══════════════════════
s = content_slide("空间智能感知 —— AI 如何理解“红线附近”", "04")
box(s, 0.5, 1.4, 12.33, 1.0, fill=DEEP, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
txt(s, 0.9, 1.58, 11.6, 0.65, [
    [("“查看潢河郝楼可采区的河道侵占情况” → ", 15, True, WHITE),
     ("一句自然语言触发完整空间分析链路", 15, True, CYAN)],
])
chain = [
    ("识别空间关键词", "郝楼可采区"),
    ("获取红线边界", "空间参考工具"),
    ("查询侵占图斑", "PostGIS 查询"),
    ("地图加载展示", "map_tool"),
    ("生成分析报告", "report_tool"),
]
x = 0.5
for i, (t, d) in enumerate(chain):
    w = 2.3
    box(s, x, 2.75, w, 1.1, fill=LIGHT, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    txt(s, x + 0.1, 2.9, w - 0.2, 0.4, [[(t, 12.5, True, DEEP)]], align=PP_ALIGN.CENTER)
    txt(s, x + 0.1, 3.3, w - 0.2, 0.4, [[(d, 10, False, GRAY)]], align=PP_ALIGN.CENTER)
    if i < 4:
        arrow(s, x + w + 0.03, 3.18, 0.28, 0.24, color=CYAN)
    x += w + 0.31
tools = [
    ("spatial_reference_tool", "河道红线 / 采区边界查询，“红线附近 / 范围内”空间语义自动识别并注入后续查询"),
    ("spatial_processing_tool", "坐标投影转换、XY 交换、CGCS2000 带号处理、GeoJSON 矢量生成与地图加载"),
    ("qgis_workflows + gis_tool_router", "QGIS MCP Docker 引擎：缓冲区 / 裁剪 / 叠加 / 相交 / 面积计算 / 分区统计 recipe 工作流"),
]
for i, (t, d) in enumerate(tools):
    y = 4.2 + i * 0.92
    box(s, 0.5, y, 12.33, 0.8, fill=PALE, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
    chip(s, 0.72, y + 0.16, 4.3, 0.48, t, MID, size=10.5)
    txt(s, 5.2, y + 0.09, 7.45, 0.62, [[(d, 11, False, INK)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)

# ══════════════════════ 第 14 页 · 应用场景 ══════════════════════
s = content_slide("应用场景", "06")
scenes = [
    ("河道侵占监管", ["遥感影像自动识别违章建筑 / 非法采砂 / 侵占水域", "自动生成检测报告", "试点：潢河郝楼、白露河黄堰采区 · 数十平方公里影像"]),
    ("采砂业务管理", ["采区边界管理、采砂许可查询、开采量统计", "“红线内是否存在违规开采”复杂空间查询", "采砂监测数据 NL2SQL 分析"]),
    ("工程报表与决策辅助", ["自动生成含地图截图、图表、结论的综合报告", "周报 / 月报周期性报表快速生成", "大幅减轻基层文档编写负担"]),
    ("知识问答与培训支持", ["水利规范 / 政策问答并引用原文条款", "新员工“随身导师”", "数据查询无结果自动回退知识库"]),
]
for i, (t, lines) in enumerate(scenes):
    x = 0.5 + (i % 2) * 6.25; y = 1.5 + (i // 2) * 2.6
    card(s, x, y, 6.0, 2.4, t, lines, accent=DEEP if i % 2 == 0 else CYAN, title_size=16, body_size=12)
box(s, 0.5, 6.6, 12.33, 0.0, fill=None)

# ══════════════════════ 第 15 页 · 项目规模 ══════════════════════
s = content_slide("项目规模与技术指标", "06")
metrics = [
    ("~24,000", "核心代码行\nPython 13,700 + React 9,600"),
    ("88+", "API 端点\n对话 / 会话 / GIS / 切片 / 知识库"),
    ("15+", "核心工具模块\n地图 / 数据库 / 检测 / 知识库"),
    ("150+", "知识库文档\n水利标准与法规"),
    ("5", "领域 Agent\n地图 / 数据 / 知识 / 报告 / 兜底"),
    ("2", "地图引擎\nLeaflet 2D + Cesium 3D"),
]
for i, (num, d) in enumerate(metrics):
    x = 0.5 + (i % 3) * 4.15; y = 1.5 + (i // 3) * 2.5
    box(s, x, y, 3.95, 2.25, fill=PALE, line=LINE, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    txt(s, x, y + 0.25, 3.95, 0.75, [[(num, 32, True, MID)]], align=PP_ALIGN.CENTER)
    txt(s, x + 0.2, y + 1.15, 3.55, 0.9, [[(ln, 11, False, INK)] for ln in d.split("\n")], align=PP_ALIGN.CENTER, line_spacing=1.15)
box(s, 0.5, 6.6, 12.33, 0.0, fill=None)
txt(s, 0.5, 6.75, 12.3, 0.4, [[("数据格式：GeoJSON · Shapefile · GeoTIFF · MBTiles · 3D Tiles · MVT    |    pytest 覆盖 Harness / RunEngine / ContextManager / RulesGateway", 10.5, False, GRAY)]])

# ══════════════════════ 第 16 页 · 总结与展望 ══════════════════════
s = content_slide("总结与未来展望", "06")
txt(s, 0.5, 1.35, 12.3, 0.45, [[("已实现 —— 完整能力闭环", 15, True, DEEP)]])
box(s, 0.5, 1.9, 12.33, 0.9, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.15)
txt(s, 0.9, 2.06, 11.6, 0.6, [[("对话式地图操作 · 遥感侵占识别 · 智能报告 · 知识库问答 · 空间分析 · 数据可视化 · 切片管理", 14, True, DEEP)]], anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.5, 3.1, 12.3, 0.45, [[("未来方向", 15, True, DEEP)]])
future = [
    ("深化变化检测", "多时相遥感月度级河道侵占自动巡查"),
    ("视频流 AI 分析", "无人机实时视频流支撑应急响应"),
    ("业务闭环", "打通政务系统数据接口\n发现 → 执法处置"),
]
for i, (t, d) in enumerate(future):
    x = 0.5 + i * 4.15
    card(s, x, 3.6, 3.95, 1.7, t, d.split("\n"), accent=AMBER, title_size=15, body_size=11.5)
box(s, 0.5, 5.75, 12.33, 1.0, fill=DEEP, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
txt(s, 0.9, 5.95, 11.6, 0.6, [
    [("核心价值：", 15, True, CYAN), ("用 AI 技术降低专业工具使用门槛，让科技真正服务于基层水利管理，为“智慧水利”与“数字乡村”建设贡献力量", 15, True, WHITE)],
], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

# ══════════════════════ 保存 ══════════════════════
prs.core_properties.title = "豫水智能一张图 - 项目总结"
prs.core_properties.author = "Map Assistant v1"
out = "/home/server/python/map_assistant_v1/docs/豫水智能一张图-项目总结.pptx"
prs.save(out)
print("saved:", out, "slides:", len(prs.slides._sldIdLst))
